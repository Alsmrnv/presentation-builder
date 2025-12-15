import hashlib
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.dml.color import RGBColor


def generate_color_scheme(presentation_name: str):
    """
    Генерирует цветовую схему на основе хеширования названия презентации
    """
    hash_input = presentation_name.encode('utf-8')
    hash_obj = hashlib.md5(hash_input)
    hex_dig = hash_obj.hexdigest()
    
    r = int(hex_dig[0:2], 16) % 256
    g = int(hex_dig[2:4], 16) % 256
    b = int(hex_dig[4:6], 16) % 256
    
    bg_factor = 0.15
    background_r = min(255, int(r * bg_factor + 240 * (1 - bg_factor)))
    background_g = min(255, int(g * bg_factor + 240 * (1 - bg_factor)))
    background_b = min(255, int(b * bg_factor + 240 * (1 - bg_factor)))
    background_color = RGBColor(background_r, background_g, background_b)
    
    pastel_factor = 0.4
    pastel_r = min(255, int(r * pastel_factor + 220 * (1 - pastel_factor)))
    pastel_g = min(255, int(g * pastel_factor + 220 * (1 - pastel_factor)))
    pastel_b = min(255, int(b * pastel_factor + 220 * (1 - pastel_factor)))
    primary_color = RGBColor(pastel_r, pastel_g, pastel_b)
    
    accent_saturation = 0.85
    accent_r = min(255, int(r * accent_saturation + 80 * (1 - accent_saturation)))
    accent_g = min(255, int(g * accent_saturation + 80 * (1 - accent_saturation)))
    accent_b = min(255, int(b * accent_saturation + 80 * (1 - accent_saturation)))
    accent_color = RGBColor(accent_r, accent_g, accent_b)
    
    text_color = RGBColor(40, 40, 40)
    
    return {
        'primary': primary_color,
        'background': background_color,
        'accent': accent_color,
        'text': text_color
    }

def apply_slide_styling(slide, color_scheme):
    """
    Применяет стили к слайду на основе цветовой схемы
    """
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color_scheme['background']
    
    try:
        slide_width = slide.width
        slide_height = slide.height
        
        slide_hash = hashlib.md5(str(slide.slide_id).encode()).hexdigest()
        angle_hex = int(slide_hash[:2], 16)
        division_angle = (angle_hex % 60) - 30
        
        primary_color = color_scheme['primary']
        r, g, b = primary_color.rgb[0], primary_color.rgb[1], primary_color.rgb[2]
        
        bg_r, bg_g, bg_b = color_scheme['background'].rgb[0], color_scheme['background'].rgb[1], color_scheme['background'].rgb[2]
        
        frag1_r = min(255, int(r * 0.05 + bg_r * 0.95))
        frag1_g = min(255, int(g * 0.05 + bg_g * 0.95))
        frag1_b = min(255, int(b * 0.05 + bg_b * 0.95))
        fragment1_color = RGBColor(frag1_r, frag1_g, frag1_b)
        
        frag2_r = min(255, int(r * 0.12 + bg_r * 0.88))
        frag2_g = min(255, int(g * 0.12 + bg_g * 0.88))
        frag2_b = min(255, int(b * 0.12 + bg_b * 0.88))
        fragment2_color = RGBColor(frag2_r, frag2_g, frag2_b)
        
        frag3_r = min(255, int(r * 0.2 + bg_r * 0.8))
        frag3_g = min(255, int(g * 0.2 + bg_g * 0.8))
        frag3_b = min(255, int(b * 0.2 + bg_b * 0.8))
        fragment3_color = RGBColor(frag3_r, frag3_g, frag3_b)
        
        diagonal_shape1 = slide.shapes.add_shape(
            1,
            Inches(-6),
            Inches(0),
            Inches(35),
            Inches(3)
        )
        diagonal_shape1.fill.solid()
        diagonal_shape1.fill.fore_color.rgb = fragment1_color
        diagonal_shape1.line.color.rgb = fragment1_color
        diagonal_shape1.line.width = Pt(0)
        diagonal_shape1.rotation = division_angle
        
        diagonal_shape2 = slide.shapes.add_shape(
            1,
            Inches(-6),
            Inches(2.5),
            Inches(35),
            Inches(2.5)
        )
        diagonal_shape2.fill.solid()
        diagonal_shape2.fill.fore_color.rgb = fragment2_color
        diagonal_shape2.line.color.rgb = fragment2_color
        diagonal_shape2.line.width = Pt(0)
        diagonal_shape2.rotation = -division_angle
        
        diagonal_shape3 = slide.shapes.add_shape(
            1,
            Inches(-6),
            Inches(4.5),
            Inches(35),
            Inches(2)
        )
        diagonal_shape3.fill.solid()
        diagonal_shape3.fill.fore_color.rgb = fragment3_color
        diagonal_shape3.line.color.rgb = fragment3_color
        diagonal_shape3.line.width = Pt(0)
        diagonal_shape3.rotation = division_angle + 10
        
        accent_bar = slide.shapes.add_shape(
            1,
            Inches(0),
            Inches(7.4),
            Inches(13.33),
            Inches(0.1)
        )
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = color_scheme['accent']
        accent_bar.line.color.rgb = color_scheme['accent']
        accent_bar.line.width = Pt(0)
        
    except Exception as e:
        print(f"Ошибка применения стиля: {e}")


def add_icon_to_slide(slide, icon_type: str, position_x: float, position_y: float, size: float = 0.5):
    """
    Добавляет иконку на слайд (в будущем будет использовать реальные иконки)
    Пока что создает геометрическую фигуру, имитирующую иконку
    """
    try:
        from pptx.enum.shapes import MSO_SHAPE
        
        shapes_map = {
            'chart': MSO_SHAPE.CHART_PLUS,
            'info': MSO_SHAPE.INFO,
            'lightning': MSO_SHAPE.BOLT,
            'star': MSO_SHAPE.STAR_4_POINT,
            'arrow': MSO_SHAPE.LEFT_ARROW,
            'check': MSO_SHAPE.CHECKBOX_X,
        }
        
        shape_type = shapes_map.get(icon_type, MSO_SHAPE.OVAL)
        
        icon_shape = slide.shapes.add_shape(
            shape_type,
            Inches(position_x),
            Inches(position_y),
            Inches(size),
            Inches(size)
        )
        
        icon_shape.fill.solid()
        icon_shape.fill.fore_color.rgb = slide.color_scheme['accent']
        icon_shape.line.color.rgb = slide.color_scheme['primary']
        icon_shape.line.width = Pt(0.5)
        
        return icon_shape
    except:
        return None

def build_presentation(slides, output_pptx: Path, presentation_name: str = "Presentation"):
    """
    Создание презентации по слайдам с поддержкой визуализаций и стилизацией
    """
    prs = Presentation()
    
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    color_scheme = generate_color_scheme(presentation_name)

    for slide_data in slides:
        has_visualization = slide_data.get("visualization", {}).get("needed", False)
        
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        
        apply_slide_styling(slide, color_scheme)
        
        slide_width = prs.slide_width
        slide_height = prs.slide_height

        if has_visualization:
            vis_data = slide_data.get("visualization", {})
            vis_type = vis_data.get("type", "")
            
            aspect_ratios = {
                "line": 50/350,
                "bar": 500/350,
                "pie": 450/350,
                "scatter": 550/350,
                "histogram": 550/350,
                "table": 700/400
            }
            
            vis_aspect_ratio = aspect_ratios.get(vis_type, 1.0)
            
            if vis_aspect_ratio >= 1.3:
                _add_slide_with_horizontal_visualization(slide, slide_data, slide_width, slide_height, color_scheme)
            else:
                _add_slide_with_vertical_visualization(slide, slide_data, slide_width, slide_height, color_scheme)
        else:
            _add_slide_full_text(slide, slide_data, slide_width, slide_height, color_scheme)

    prs.save(output_pptx.as_posix())
    print(f"✅ Презентация сохранена: {output_pptx}")

def _add_slide_full_text(slide, slide_data, slide_width, slide_height, color_scheme):
    """Добавляет слайд с полноразмерным текстом (когда нет визуализации)"""
    margin_left = Inches(0.8)
    margin_right = Inches(0.8)
    top_margin = Inches(0.6)
    bottom_margin = Inches(0.8)
    between = Inches(0.3)
    title_height = Inches(0.9)

    title_left = margin_left
    title_top = top_margin
    title_width = slide_width - margin_left - margin_right
    body_left = margin_left
    body_top = title_top + title_height + between
    body_width = title_width
    body_height = slide_height - body_top - bottom_margin

    title_shape = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    title_tf = title_shape.text_frame
    title_tf.clear()
    title_tf.text = slide_data["title"]
    title_tf.word_wrap = True
    title_tf.auto_size = MSO_AUTO_SIZE.NONE
    
    for paragraph in title_tf.paragraphs:
        paragraph.alignment = 1
        for run in paragraph.runs:
            run.font.color.rgb = color_scheme['accent']
            run.font.name = 'Calibri'
            run.font.bold = True
            run.font.size = Pt(32)

    body_shape = slide.shapes.add_textbox(body_left, body_top, body_width, body_height)
    tf = body_shape.text_frame
    tf.clear()
    tf.text = slide_data["description"]
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    text_len = len(slide_data["description"])
    if text_len > 800:
        base_size = 20
    elif text_len > 500:
        base_size = 22
    elif text_len > 300:
        base_size = 24
    else:
        base_size = 26

    for paragraph in tf.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(base_size)
            run.font.color.rgb = color_scheme['text']
            run.font.name = 'Calibri'


def _add_slide_with_vertical_visualization(slide, slide_data, slide_width, slide_height, color_scheme):
    """Добавляет слайд с вертикальной визуализацией (рядом с текстом)"""
    margin = Inches(0.5)
    between = Inches(0.3)
    title_height = Inches(0.7)
    
    vis_data = slide_data.get("visualization", {})
    
    text_width = slide_width * 0.5 - margin
    image_width = slide_width * 0.5 - margin
    text_left = margin
    image_left = text_left + text_width + between
    
    title_left = margin
    title_top = margin
    title_width = slide_width - 2 * margin
    
    text_top = title_top + title_height + between
    text_height = slide_height - text_top - margin
    image_top = text_top
    image_height = text_height
    
    title_shape = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    title_tf = title_shape.text_frame
    title_tf.clear()
    title_tf.text = slide_data["title"]
    title_tf.word_wrap = True
    title_tf.auto_size = MSO_AUTO_SIZE.NONE
    
    for paragraph in title_tf.paragraphs:
        paragraph.alignment = 1
        for run in paragraph.runs:
            run.font.color.rgb = color_scheme['accent']
            run.font.name = 'Calibri'
            run.font.bold = True
            run.font.size = Pt(28)

    text_shape = slide.shapes.add_textbox(text_left, text_top, text_width, text_height)
    text_tf = text_shape.text_frame
    text_tf.clear()
    text_tf.text = slide_data["description"]
    text_tf.word_wrap = True
    text_tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    
    text_len = len(slide_data["description"])
    if text_len > 800:
        base_size = 16
    elif text_len > 500:
        base_size = 18
    elif text_len > 300:
        base_size = 20
    else:
        base_size = 22
    
    for paragraph in text_tf.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(base_size)
            run.font.color.rgb = color_scheme['text']
            run.font.name = 'Calibri'
    
    image_path = vis_data.get("image_path")
    
    if image_path and Path(image_path).exists():
        try:
            from PIL import Image
            img = Image.open(image_path)
            img_width, img_height = img.size
            img_aspect = img_width / img_height
            
            available_aspect = image_width / image_height
            
            if img_aspect > available_aspect:
                final_width = image_width
                final_height = image_width / img_aspect
            else:
                final_height = image_height
                final_width = image_height * img_aspect
            
            img_left = image_left + (image_width - final_width) / 2
            img_top = image_top + (image_height - final_height) / 2
            
            slide.shapes.add_picture(
                str(image_path),
                img_left,
                img_top,
                width=final_width,
                height=final_height
            )
        except Exception as e:
            print(f"Ошибка добавления картинки {image_path}: {e}")
    else:
        print(f"Картинка не найдена: {image_path}")


def _add_slide_with_horizontal_visualization(slide, slide_data, slide_width, slide_height, color_scheme):
    """Добавляет слайд с горизонтальной визуализацией (под текстом)"""
    margin = Inches(0.5)
    between = Inches(0.3)
    title_height = Inches(0.7)
    
    vis_data = slide_data.get("visualization", {})
    
    title_left = margin
    title_top = margin
    title_width = slide_width - 2 * margin
    
    text_left = margin
    text_top = title_top + title_height + between
    text_width = title_width
    
    title_shape = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    title_tf = title_shape.text_frame
    title_tf.clear()
    title_tf.text = slide_data["title"]
    title_tf.word_wrap = True
    title_tf.auto_size = MSO_AUTO_SIZE.NONE
    
    for paragraph in title_tf.paragraphs:
        paragraph.alignment = 1
        for run in paragraph.runs:
            run.font.color.rgb = color_scheme['accent']
            run.font.name = 'Calibri'
            run.font.bold = True
            run.font.size = Pt(28)

    text_shape = slide.shapes.add_textbox(text_left, text_top, text_width, Inches(1.5))
    text_tf = text_shape.text_frame
    text_tf.clear()
    text_tf.text = slide_data["description"]
    text_tf.word_wrap = True
    text_tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    
    text_len = len(slide_data["description"])
    if text_len > 800:
        base_size = 16
    elif text_len > 500:
        base_size = 18
    elif text_len > 300:
        base_size = 20
    else:
        base_size = 22

    for paragraph in text_tf.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(base_size)
            run.font.color.rgb = color_scheme['text']
            run.font.name = 'Calibri'
    
    text_shape_height = min(Inches(2.0), slide_height - title_top - title_height - between - Inches(2.0))
    image_top = text_top + text_shape_height + between
    image_left = margin
    image_width = slide_width - 2 * margin
    image_height = slide_height - image_top - margin
    
    image_path = vis_data.get("image_path")
    
    if image_path and Path(image_path).exists():
        try:
            from PIL import Image
            img = Image.open(image_path)
            img_width, img_height = img.size
            img_aspect = img_width / img_height
            
            available_aspect = image_width / image_height
            
            if img_aspect > available_aspect:
                final_width = image_width
                final_height = image_width / img_aspect
            else:
                final_height = image_height
                final_width = image_height * img_aspect
            
            img_left = image_left + (image_width - final_width) / 2
            img_top = image_top + (image_height - final_height) / 2
            
            slide.shapes.add_picture(
                str(image_path),
                img_left,
                img_top,
                width=final_width,
                height=final_height
            )
        except Exception as e:
            print(f"Ошибка добавления картинки {image_path}: {e}")
    else:
        print(f"Картинка не найдена: {image_path}")

